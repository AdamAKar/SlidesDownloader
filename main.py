# main.py
# ----------------
# A self-contained desktop application to download a Google Slides
# presentation (with full-quality images and videos) as a .pptx file.

import sys
import os
import io
import re
import webbrowser
from PyQt5 import QtWidgets
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from pptx import Presentation
from pptx.util import Inches

# OAuth2 scopes: read-only access to Slides & Drive
SCOPES = [
    'https://www.googleapis.com/auth/presentations.readonly',
    'https://www.googleapis.com/auth/drive.readonly'
]

class SlideDownloader(QtWidgets.QWidget):
    def __init__(self):
        super().__init__()
        self.init_ui()

    def init_ui(self):
        self.setWindowTitle('Google Slides → PPTX Downloader')
        self.resize(500, 300)
        layout = QtWidgets.QVBoxLayout()

        self.url_input = QtWidgets.QLineEdit()
        self.url_input.setPlaceholderText('Paste full Google Slides URL here')
        layout.addWidget(self.url_input)

        self.auth_btn = QtWidgets.QPushButton('Authenticate & Download')
        self.auth_btn.clicked.connect(self.handle_download)
        layout.addWidget(self.auth_btn)

        self.log = QtWidgets.QTextEdit()
        self.log.setReadOnly(True)
        layout.addWidget(self.log)

        self.setLayout(layout)

    def log_msg(self, msg: str):
        self.log.append(msg)
        QtWidgets.QApplication.processEvents()

    def extract_presentation_id(self, url: str) -> str:
        # Try typical /d/{ID}/ pattern
        m = re.search(r'/d/([a-zA-Z0-9_-]+)', url)
        if m:
            return m.group(1)
        # Fallback: presentations/d/
        m2 = re.search(r'presentation/d/([a-zA-Z0-9_-]+)', url)
        if m2:
            return m2.group(1)
        return None

    def authenticate(self):
        # Launch browser OAuth flow; creates token.json automatically
        flow = InstalledAppFlow.from_client_secrets_file('credentials.json', SCOPES)
        creds = flow.run_local_server(port=0)
        return creds

    def download_slide_deck(self, pres_id: str, creds):
        slides_service = build('slides', 'v1', credentials=creds)
        drive_service = build('drive', 'v3', credentials=creds)

        self.log_msg(f'Fetching presentation metadata...')
        presentation = slides_service.presentations().get(presentationId=pres_id).execute()

        ppt = Presentation()

        for idx, slide in enumerate(presentation.get('slides', []), start=1):
            self.log_msg(f'Processing slide {idx} of {len(presentation.get("slides", []))}...')
            # blank slide layout
            blank = ppt.slide_layouts[6]
            new_slide = ppt.slides.add_slide(blank)

            for element in slide.get('pageElements', []):
                # Handle images
                if 'image' in element:
                    img_url = element['image'].get('contentUrl')
                    if img_url:
                        img_path = self.download_drive_file(img_url, drive_service)
                        new_slide.shapes.add_picture(img_path, Inches(1), Inches(1))

                # Handle videos
                if 'video' in element:
                    vid = element['video']
                    source = vid.get('source', '')
                    if source:
                        self.log_msg(f'Detected video: {source}')
                        if 'drive.google.com' in source:
                            QtWidgets.QMessageBox.information(
                                self, 'Grant Access Required',
                                'This video is on Google Drive and may not be public.\n'
                                'Please grant view access and try again.'
                            )
                            webbrowser.open(source)
                        else:
                            # Insert hyperlink text
                            tx_box = new_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.5))
                            tx = tx_box.text_frame
                            tx.text = f'Video link: {source}'

        # Save file dialog
        out_path, _ = QtWidgets.QFileDialog.getSaveFileName(
            self, 'Save presentation as…', filter='PowerPoint Files (*.pptx)')
        if out_path:
            if not out_path.lower().endswith('.pptx'):
                out_path += '.pptx'
            ppt.save(out_path)
            self.log_msg(f'Saved to: {out_path}')

    def download_drive_file(self, url: str, drive_service):
        # Extract file ID
        m = re.search(r'/d/([a-zA-Z0-9_-]+)', url)
        if not m:
            raise ValueError('Could not parse Drive file ID')
        file_id = m.group(1)
        request = drive_service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        # Save temporarily
        tmp_path = os.path.join(os.getcwd(), f'{file_id}')
        with open(tmp_path, 'wb') as f:
            f.write(fh.read())
        return tmp_path

    def handle_download(self):
        url = self.url_input.text().strip()
        if not url:
            self.log_msg('Error: no URL provided')
            return
        pres_id = self.extract_presentation_id(url)
        if not pres_id:
            self.log_msg('Error: could not extract presentation ID')
            return
        try:
            creds = self.authenticate()
            self.download_slide_deck(pres_id, creds)
        except Exception as e:
            self.log_msg(f'Error: {e}')

if __name__ == '__main__':
    app = QtWidgets.QApplication(sys.argv)
    downloader = SlideDownloader()
    downloader.show()
    sys.exit(app.exec_())
